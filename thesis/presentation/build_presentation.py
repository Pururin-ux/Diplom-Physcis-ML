from __future__ import annotations

import zipfile
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "thesis" / "presentation"
ASSETS = ROOT / "reports" / "assets"
GEN = OUT_DIR / "generated_assets"
PPTX = OUT_DIR / "diploma_defense_presentation.pptx"
OLD_PDF = OUT_DIR / "diploma_defense_presentation.pdf"
OLD_TEX = OUT_DIR / "diploma_defense_presentation.tex"

TEMPLATE_CANDIDATES = [
    Path(r"C:\Users\lalad\Downloads\Telegram Desktop\Презентация (Сиротюк) (2).pptx"),
    Path(r"C:\Users\lalad\Downloads\Telegram Desktop\Презентация (Сиротюк).pptx"),
]

W, H = 13.333333, 7.5
BLUE = RGBColor(0x00, 0x20, 0x60)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
GRID = RGBColor(0xC9, 0xD3, 0xE0)


def inch(value: float):
    return Inches(value)


def template_path() -> Path | None:
    for candidate in TEMPLATE_CANDIDATES:
        if candidate.exists():
            return candidate
    return None


def ensure_template_logos(tpl: Path | None) -> None:
    """Extract only the template logos needed on the title slide."""
    media_dir = OUT_DIR / "template_media"
    media_dir.mkdir(parents=True, exist_ok=True)
    if tpl is None or not tpl.exists():
        return
    wanted = {
        "ppt/media/image2.png": media_dir / "image2.png",
        "ppt/media/image5.png": media_dir / "image5.png",
    }
    with zipfile.ZipFile(tpl, "r") as zf:
        for src, dst in wanted.items():
            if src in zf.namelist() and not dst.exists():
                dst.write_bytes(zf.read(src))


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # python-pptx has no public delete API.
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def set_run_font(run, size: float, bold: bool = False, color: RGBColor = BLACK) -> None:
    run.font.name = "Times New Roman"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 20,
    bold: bool = False,
    color: RGBColor = BLACK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.05)
    tf.margin_right = inch(0.05)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_run_font(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title: str, num: int) -> None:
    add_text(slide, title, 0.333, 0.045, 11.85, 0.78, size=32, bold=False)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(12.396), inch(0.0), inch(0.937), inch(0.937))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(num)
    set_run_font(r, size=22, bold=True, color=WHITE)


def add_box(slide, x: float, y: float, w: float, h: float, fill=WHITE, line=BLUE, width=1.1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, size: float = 21) -> None:
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.05)
    tf.margin_right = inch(0.05)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.text = bullet
        p.font.name = "Times New Roman"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(5)


def add_caption(slide, text: str, x: float, y: float, w: float, h: float = 0.45, size: float = 16) -> None:
    add_text(slide, text, x, y, w, h, size=size, align=PP_ALIGN.CENTER)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), inch(x), inch(y), width=inch(w), height=inch(h))


def add_key_line(slide, text: str) -> None:
    add_box(slide, 0.35, 6.78, 12.60, 0.52, fill=LIGHT, line=BLUE, width=1.0)
    add_text(slide, text, 0.55, 6.85, 12.20, 0.35, size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def generate_support_assets() -> None:
    GEN.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update(
        {
            "font.family": "Times New Roman",
            "mathtext.fontset": "dejavuserif",
            "axes.labelsize": 16,
            "xtick.labelsize": 13,
            "ytick.labelsize": 13,
        }
    )

    fig, ax = plt.subplots(figsize=(6.6, 4.3), dpi=220)
    ax.set_aspect("equal")
    a, b, n = 6.6, 4.1, 3.0
    xs = np.arange(-8, 9)
    ys = np.arange(-6, 7)
    xx, yy = np.meshgrid(xs, ys)
    inside = (np.abs(xx / a) ** n + np.abs(yy / b) ** n) <= 1
    ax.scatter(xx[~inside], yy[~inside], s=24, color="#E4E9F0", edgecolors="#8BA5C4", linewidths=0.5)
    ax.scatter(xx[inside], yy[inside], s=46, color="#00458F", edgecolors="white", linewidths=0.5)
    t = np.linspace(0, 2 * np.pi, 600)
    bx = np.sign(np.cos(t)) * np.abs(np.cos(t)) ** (2 / n) * a
    by = np.sign(np.sin(t)) * np.abs(np.sin(t)) ** (2 / n) * b
    ax.plot(bx, by, color="#111111", lw=2.0)
    ax.set_xlim(-8.8, 8.8)
    ax.set_ylim(-6.4, 6.4)
    ax.axis("off")
    ax.text(-8.4, 5.55, "узлы внутри области", color="#00458F", fontsize=15, weight="bold")
    ax.text(-8.4, -5.75, "открытая граница", color="#111111", fontsize=15)
    fig.tight_layout(pad=0.05)
    fig.savefig(GEN / "model_lattice_superellipse.png", bbox_inches="tight")
    plt.close(fig)

    fig, axes = plt.subplots(1, 2, figsize=(8.8, 4.0), dpi=220)
    for ax, title, mode in zip(
        axes,
        ["LOAO: исключён один размер $a$", "LOARO: исключено одно отношение сторон"],
        ["col", "row"],
    ):
        data = np.zeros((7, 5))
        if mode == "col":
            data[:, 2] = 1
        else:
            data[4, :] = 1
        ax.imshow(data, cmap=plt.matplotlib.colors.ListedColormap(["#E8EDF4", "#00458F"]), vmin=0, vmax=1)
        ax.set_xticks(range(5), ["24", "27", "30", "33", "36"])
        ax.set_yticks(range(7), ["0.67", "0.72", "0.78", "0.83", "0.89", "0.94", "1.0"])
        ax.set_xlabel("$a$")
        ax.set_ylabel(r"$r_{\mathrm{AR}}$")
        ax.set_title(title, fontsize=16, color="#002060", pad=7)
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.tick_params(length=0)
        ax.set_xticks(np.arange(-0.5, 5, 1), minor=True)
        ax.set_yticks(np.arange(-0.5, 7, 1), minor=True)
        ax.grid(which="minor", color="white", linewidth=2)
    fig.tight_layout(pad=0.8, w_pad=1.8)
    fig.savefig(GEN / "loao_loaro_grid.png", bbox_inches="tight")
    plt.close(fig)

    fig, ax = plt.subplots(figsize=(6.5, 3.5), dpi=220)
    labels = ["LOAO", "LOARO", "всего"]
    vals = [3 / 8 * 100, 7 / 8 * 100, 10 / 16 * 100]
    counts = ["3/8", "7/8", "10/16"]
    colors = ["#A6A6A6", "#00458F", "#00458F"]
    bars = ax.bar(labels, vals, color=colors, width=0.55)
    ax.set_ylim(0, 100)
    ax.set_ylabel("успешные ячейки, %")
    ax.grid(axis="y", alpha=0.25)
    for b, c in zip(bars, counts):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 3, c, ha="center", fontsize=16, weight="bold")
    ax.set_title("Критерий выполнен не во всех схемах проверки", color="#002060", fontsize=17)
    fig.tight_layout()
    fig.savefig(GEN / "mlp_success_summary_template.png", bbox_inches="tight")
    plt.close(fig)


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, font_size: float = 17) -> None:
    shape = slide.shapes.add_table(len(rows), len(rows[0]), inch(x), inch(y), inch(w), inch(h))
    table = shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = inch(0.04)
            cell.margin_right = inch(0.04)
            cell.margin_top = inch(0.02)
            cell.margin_bottom = inch(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE9, 0xEF, 0xF7) if r == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run_font(run, size=font_size, bold=(r == 0), color=BLACK)


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logos = OUT_DIR / "template_media"
    if (logos / "image5.png").exists():
        slide.shapes.add_picture(str(logos / "image5.png"), inch(4.55), inch(0.75), width=inch(2.85))
    if (logos / "image2.png").exists():
        slide.shapes.add_picture(str(logos / "image2.png"), inch(7.58), inch(0.58), width=inch(3.55))
    add_text(slide, "ДИПЛОМНАЯ РАБОТА", 0.38, 1.75, 12.58, 0.45, size=22, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "Моделирование энергетического спектра\nи волновых функций в квантовых точках\nсложной геометрии",
        0.38,
        2.25,
        12.58,
        1.70,
        size=28,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Модельные суперэллиптические квантовые точки\nна квадратной решётке метода сильной связи",
        0.80,
        4.05,
        11.70,
        0.80,
        size=22,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Докладчик:\nАльвинский Александр Андреевич, студент 5 курса\n\nНаучный руководитель:\nА. В. Ларькин, канд. физ.-мат. наук",
        0.55,
        5.05,
        10.50,
        1.75,
        size=21,
    )


def slide_goal(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Цель и задачи работы", 2)
    add_box(slide, 0.45, 1.18, 12.40, 1.15, fill=LIGHT, line=BLUE)
    add_text(
        slide,
        "Форма и размер модельной квантовой точки управляют\nнизкоэнергетическим спектром за счёт квантового ограничения.",
        0.75,
        1.36,
        11.80,
        0.75,
        size=24,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Цель",
        0.55,
        2.75,
        5.45,
        0.38,
        size=22,
        bold=True,
        color=BLUE,
    )
    add_box(slide, 0.55, 3.12, 5.55, 0.03, fill=BLUE, line=BLUE)
    add_text(
        slide,
        "Смоделировать энергетический спектр и волновые функции квантовых точек сложной геометрии в рамках метода сильной связи.",
        0.55,
        3.35,
        5.65,
        1.65,
        size=22,
    )
    add_text(slide, "Задачи", 6.75, 2.75, 5.70, 0.38, size=22, bold=True, color=BLUE)
    add_box(slide, 6.75, 3.12, 5.70, 0.03, fill=BLUE, line=BLUE)
    add_bullets(
        slide,
        [
            "построить модель сильной связи;",
            "проверить физическую согласованность расчётов;",
            "сравнить Ridge и MLP при LOAO/LOARO.",
        ],
        6.75,
        3.32,
        5.95,
        1.65,
        size=22,
    )
    add_key_line(slide, "Сначала физическая проверка, затем суррогатная аппроксимация.")


def slide_relevance(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Актуальность работы", 3)
    add_bullets(
        slide,
        [
            "Квантовое ограничение делает спектр конечной наноструктуры чувствительным к размеру и форме.",
            "Модельная решётка позволяет отделить влияние геометрии от материал-специфических параметров.",
            "Суррогатная модель нужна только как быстрая аппроксимация результатов прямого расчёта.",
        ],
        0.55,
        1.25,
        6.0,
        2.6,
        size=22,
    )
    add_box(slide, 6.85, 1.20, 5.95, 4.70, fill=LIGHT, line=GRID, width=0.8)
    add_text(
        slide,
        "Физическая цепочка",
        7.10,
        1.45,
        5.45,
        0.35,
        size=23,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    for k, txt in enumerate(["размер и форма", "квантовое ограничение", "низкие уровни", "суррогатная проверка"]):
        y = 2.05 + k * 0.78
        add_box(slide, 7.45, y, 4.75, 0.48, fill=WHITE, line=BLUE)
        add_text(slide, txt, 7.55, y + 0.08, 4.55, 0.25, size=19, bold=(k == 2), align=PP_ALIGN.CENTER)
        if k < 3:
            add_text(slide, "↓", 9.72, y + 0.47, 0.4, 0.22, size=18, color=BLUE, align=PP_ALIGN.CENTER)
    add_key_line(slide, "Суррогат не заменяет прямой квантовый расчёт.")


def slide_model(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Модель метода сильной связи", 4)
    add_text(
        slide,
        "H = Σ εᵢ |i⟩⟨i| + Σ tᵢⱼ (|i⟩⟨j| + |j⟩⟨i|)\n\nεᵢ = 0,    tᵢⱼ = −1\n\nEₖᵢₙ = E₀ + 4",
        0.55,
        1.35,
        5.25,
        2.55,
        size=24,
        align=PP_ALIGN.CENTER,
    )
    add_box(slide, 0.85, 4.35, 4.65, 1.15, fill=LIGHT, line=BLUE)
    add_text(
        slide,
        "Конечная область решётки\nрассматривается как модельная\nквантовая точка.",
        1.00,
        4.55,
        4.35,
        0.75,
        size=20,
        align=PP_ALIGN.CENTER,
    )
    add_image_fit(slide, GEN / "model_lattice_superellipse.png", 6.20, 1.25, 6.45, 4.85)
    add_caption(slide, "Рисунок 1 – дискретная область на квадратной решётке", 6.25, 6.12, 6.35, 0.35, size=15)
    add_key_line(slide, "Спектр определяется размером, отношением сторон и формой границы.")


def slide_dataset(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Расчётная выборка", 5)
    add_box(slide, 0.55, 1.25, 5.20, 4.95, fill=LIGHT, line=BLUE)
    add_text(slide, "140 геометрий", 0.90, 1.62, 4.55, 0.55, size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    rows = [
        ["параметр", "значения"],
        ["n", "1.2; 2.0; 3.0; 4.0"],
        ["a", "24; 27; 30; 33; 36"],
        ["r_AR", "0.67; 0.72; … ; 1.0"],
        ["b", "a r_AR"],
    ]
    add_table(slide, rows, 0.90, 2.45, 4.50, 2.55, font_size=16)
    add_text(
        slide,
        "n рассматривается как фиксированный дискретный класс,\nа не как непрерывный вход модели.",
        0.85,
        5.35,
        4.60,
        0.50,
        size=17,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "|x/a|ⁿ + |y/b|ⁿ ≤ 1,      b = a r_AR",
        6.20,
        1.55,
        6.30,
        0.60,
        size=25,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_image_fit(slide, ASSETS / "nsites_area_ratio_by_n.png", 6.15, 2.25, 6.20, 3.60)
    add_caption(slide, "Рисунок 2 – контроль дискретной площади области", 6.15, 5.95, 6.20, 0.35, size=15)
    add_key_line(slide, "Выборка покрывает проверенный диапазон параметров, но не задаёт универсальную модель формы.")


def slide_energy(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Энергетическая шкала", 6)
    add_text(
        slide,
        "E(kₓ,kᵧ) = −2 cos kₓ − 2 cos kᵧ\n\nE_min = −4\n\nE(kₓ,kᵧ)+4 ≈ kₓ²+kᵧ²\n\nE_kin = E₀ + 4",
        0.55,
        1.25,
        5.75,
        3.45,
        size=24,
        align=PP_ALIGN.CENTER,
    )
    add_bullets(
        slide,
        [
            "низкие уровни отсчитываются от дна зоны;",
            "при увеличении размера ослабляется квантовое ограничение;",
            "это мотивирует масштаб 1/a².",
        ],
        0.65,
        5.05,
        5.70,
        1.05,
        size=20,
    )
    add_image_fit(slide, ASSETS / "e0_kin_a2_scaling_by_n.png", 6.55, 1.22, 6.20, 4.65)
    add_caption(slide, "Рисунок 3 – проверка масштаба (E₀+4)a²", 6.60, 6.00, 6.10, 0.35, size=15)
    add_key_line(slide, "Проверяется не полный континуальный предел, а согласованность низкоэнергетического масштаба.")


def slide_physics_checks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Физические проверки", 7)
    add_image_fit(slide, ASSETS / "ar_scaling_relative_deviation.png", 0.55, 1.15, 6.15, 4.95)
    add_caption(slide, "Рисунок 4 – отклонение (E₀+4)a² при фиксированных n и r_AR", 0.55, 6.15, 6.15, 0.35, size=14)
    add_image_fit(slide, ASSETS / "circle_bessel_e0_check.png", 7.05, 1.15, 5.65, 3.30)
    add_caption(slide, "Рисунок 5 – круговая проверка Бесселя", 7.05, 4.50, 5.65, 0.35, size=14)
    add_bullets(
        slide,
        [
            "круговой случай: ошибка ≈ 2.03%;",
            "максимальное отклонение по фиксированной форме ≈ 2.12%;",
            "dE₂ оставлен только как диагностическая величина.",
        ],
        7.05,
        5.05,
        5.85,
        1.25,
        size=18,
    )
    add_key_line(slide, "Физические проверки поддерживают выбранную область параметров.")


def slide_targets(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Спектральные характеристики", 8)
    add_box(slide, 0.65, 1.30, 3.75, 4.55, fill=LIGHT, line=BLUE)
    add_text(slide, "E₀", 0.95, 1.65, 3.15, 0.55, size=34, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "нижний уровень\nквантового ограничения", 0.90, 2.40, 3.25, 1.00, size=22, align=PP_ALIGN.CENTER)
    add_box(slide, 4.80, 1.30, 3.75, 4.55, fill=LIGHT, line=BLUE)
    add_text(slide, "dE₁ = E₁ − E₀", 5.05, 1.65, 3.25, 0.55, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "зазор до первого\nвозбуждённого состояния", 5.05, 2.40, 3.25, 1.00, size=22, align=PP_ALIGN.CENTER)
    add_box(slide, 8.95, 1.30, 3.75, 4.55, fill=LIGHT, line=BLUE)
    add_text(slide, "dE₂", 9.25, 1.65, 3.15, 0.55, size=34, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "диагностика\nвырождения и порядка\nуровней", 9.20, 2.35, 3.25, 1.30, size=22, align=PP_ALIGN.CENTER)
    add_image_fit(slide, ASSETS / "de2_near_degeneracy_n2_ar1.png", 9.50, 4.10, 2.65, 1.30)
    add_key_line(slide, "Основные целевые величины: E₀ и dE₁; dE₂ не используется как главная цель регрессии.")


def slide_surrogate(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Суррогатная аппроксимация", 9)
    add_box(slide, 0.55, 1.20, 5.85, 1.35, fill=LIGHT, line=BLUE)
    add_text(slide, "Признаки", 0.75, 1.35, 5.45, 0.30, size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "[ 1/a²,   1/(a² r_AR),   r_AR ]", 0.75, 1.82, 5.45, 0.38, size=25, align=PP_ALIGN.CENTER)
    add_box(slide, 0.55, 3.05, 2.65, 1.45, fill=WHITE, line=BLUE)
    add_text(slide, "Ridge", 0.75, 3.25, 2.25, 0.35, size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "основная\nинтерпретируемая\nмодель", 0.75, 3.70, 2.25, 0.60, size=18, align=PP_ALIGN.CENTER)
    add_box(slide, 3.75, 3.05, 2.65, 1.45, fill=WHITE, line=BLUE)
    add_text(slide, "MLP", 3.95, 3.25, 2.25, 0.35, size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "малая контрольная\nнелинейная\nмодель", 3.95, 3.70, 2.25, 0.60, size=18, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "суррогат обучается по расчётным спектрам;",
            "Kwant остаётся эталонным расчётом;",
            "коэффициенты Ridge не трактуются как отдельные физические константы.",
        ],
        0.65,
        5.05,
        5.95,
        1.15,
        size=18,
    )
    add_image_fit(slide, GEN / "loao_loaro_grid.png", 6.70, 1.20, 6.05, 4.65)
    add_caption(slide, "Рисунок 6 – структурированные схемы проверки", 6.70, 5.95, 6.05, 0.35, size=14)
    add_key_line(slide, "Случайное разбиение не использовалось из-за регулярной сетки параметров.")


def slide_mlp_result(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Сравнение Ridge и MLP", 10)
    add_image_fit(slide, ASSETS / "mlp_ablation_improvement_by_cell.png", 0.55, 1.20, 7.70, 4.80)
    add_caption(slide, "Рисунок 7 – относительное улучшение MLP по сравнению с Ridge", 0.55, 6.08, 7.70, 0.35, size=14)
    add_box(slide, 8.60, 1.25, 4.10, 4.90, fill=LIGHT, line=BLUE)
    add_text(slide, "Итог критерия", 8.85, 1.55, 3.60, 0.35, size=23, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "MLP + физические признаки: 10/16;",
            "LOAO: 3/8;",
            "LOARO: 7/8;",
            "заданный критерий не выполнен.",
        ],
        8.85,
        2.15,
        3.55,
        2.05,
        size=19,
    )
    add_image_fit(slide, GEN / "mlp_success_summary_template.png", 8.90, 4.42, 3.50, 1.45)
    add_key_line(slide, "MLP улучшает часть ячеек, но не даёт устойчивого преимущества.")


def slide_error_scale(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Физический масштаб ошибок", 11)
    rows = [
        ["модель", "E₀ от E_kin", "dE₁ от dE₁"],
        ["Ridge", "2.01–2.60%", "0.78–1.15%"],
        ["MLP + физ. признаки", "0.51–2.06%", "0.49–3.31%"],
    ]
    add_table(slide, rows, 0.55, 1.35, 6.10, 1.75, font_size=18)
    add_bullets(
        slide,
        [
            "ошибки MAE заданы в единицах |t|;",
            "для E₀ масштабом служит E_kin = E₀ + 4, а не |E₀|;",
            "для dE₁ ошибка нормируется на характерный зазор.",
        ],
        0.70,
        3.55,
        5.75,
        1.75,
        size=20,
    )
    add_image_fit(slide, ASSETS / "mlp_ablation_ridge_vs_mlp_physics_mae.png", 6.95, 1.25, 5.85, 4.70)
    add_caption(slide, "Рисунок 8 – MAE Ridge и MLP с физическими признаками", 6.95, 6.05, 5.85, 0.35, size=14)
    add_key_line(slide, "Абсолютные ошибки порядка 10⁻⁴ |t| имеют понятный физический масштаб.")


def slide_residuals(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Диагностика остатков Ridge", 12)
    add_image_fit(slide, ASSETS / "ridge_residual_dominance_by_n.png", 0.55, 1.20, 6.45, 4.70)
    add_caption(slide, "Рисунок 9 – доминирование диагностик по классам n", 0.55, 6.00, 6.45, 0.35, size=14)
    add_box(slide, 7.35, 1.25, 5.25, 4.85, fill=LIGHT, line=BLUE)
    add_text(slide, "Основные выводы", 7.65, 1.55, 4.65, 0.35, size=23, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "глобальная гипотеза об остатках не поддержана;",
            "диагностический сигнал найден только для n = 1.2: ρ ≈ −0.416, p ≈ 0.013;",
            "для n ≥ 2 простые диагностики границы не объясняют остатки систематически.",
        ],
        7.60,
        2.25,
        4.70,
        2.35,
        size=19,
    )
    add_text(
        slide,
        "Ограничения: только заданная сетка параметров; без непрерывной модели по n; без DFT/OpenMX-калибровки.",
        7.60,
        5.10,
        4.70,
        0.65,
        size=17,
        align=PP_ALIGN.CENTER,
    )
    add_key_line(slide, "Прямой расчёт Kwant остаётся эталоном; суррогат служит контролируемой аппроксимацией.")


def cleanup_stale_latex_outputs() -> None:
    for path in [OLD_PDF, OLD_TEX]:
        if path.exists():
            path.unlink()
    for suffix in [".aux", ".log", ".nav", ".out", ".snm", ".toc"]:
        path = OUT_DIR / f"diploma_defense_presentation{suffix}"
        if path.exists():
            path.unlink()
    for path in GEN.glob("latex_slide-*.png"):
        path.unlink()


def build() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    generate_support_assets()

    tpl = template_path()
    ensure_template_logos(tpl)
    if tpl is not None:
        prs = Presentation(str(tpl))
        delete_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = inch(W)
        prs.slide_height = inch(H)

    slide_title(prs)
    slide_goal(prs)
    slide_relevance(prs)
    slide_model(prs)
    slide_dataset(prs)
    slide_energy(prs)
    slide_physics_checks(prs)
    slide_targets(prs)
    slide_surrogate(prs)
    slide_mlp_result(prs)
    slide_error_scale(prs)
    slide_residuals(prs)

    prs.save(PPTX)
    cleanup_stale_latex_outputs()
    print(PPTX)
    print(f"slides: {len(prs.slides)}")
    print("pdf export: not available in this environment")


if __name__ == "__main__":
    build()
