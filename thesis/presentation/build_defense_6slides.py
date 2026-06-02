# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
import zipfile

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "thesis" / "presentation"
ASSETS = ROOT / "reports" / "assets"
GEN = OUT_DIR / "generated_assets"
PREVIEW = OUT_DIR / "preview_png"
OUT = OUT_DIR / "diploma_defense_6slides.pptx"
TEMPLATE = Path(r"C:\Users\lalad\Downloads\Telegram Desktop\Презентация (Сиротюк) (2).pptx")

BLUE = RGBColor(0x1E, 0x5A, 0xA5)
DARK_BLUE = RGBColor(0x00, 0x20, 0x60)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(0xEA, 0xF1, 0xF8)
MID_BLUE = RGBColor(0xB8, 0xCC, 0xE4)
GRAY = RGBColor(0x8D, 0x8D, 0x8D)
FONT = "Times New Roman"
SLIDE_W = 13.333333
SLIDE_H = 7.5


def inch(value: float):
    return Inches(value)


def delete_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def set_font(run, size: float, bold: bool = False, color: RGBColor = BLACK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.color.rgb = color


def text_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 24,
    bold: bool = False,
    color: RGBColor = BLACK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.03,
):
    shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = valign
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    for i, line_text in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line_text
        set_font(r, size=size, bold=bold, color=color)
    return shape


def box(slide, x, y, w, h, fill=None, line_color=None, width=0.7):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None or width == 0:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(width)
    return shape


def line(slide, x1, y1, x2, y2, color=DARK_BLUE, width=1.0):
    shape = slide.shapes.add_connector(1, inch(x1), inch(y1), inch(x2), inch(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def add_picture(slide, path: Path, x: float, y: float, w: float, h: float | None = None):
    if h is None:
        return slide.shapes.add_picture(str(path), inch(x), inch(y), width=inch(w))
    return slide.shapes.add_picture(str(path), inch(x), inch(y), width=inch(w), height=inch(h))


def ensure_template_logos() -> None:
    if not TEMPLATE.exists():
        return
    media = OUT_DIR / "template_media"
    media.mkdir(parents=True, exist_ok=True)
    wanted = {
        "ppt/media/image2.png": media / "image2.png",
        "ppt/media/image5.png": media / "image5.png",
        "ppt/media/image6.png": media / "image6.png",
    }
    with zipfile.ZipFile(TEMPLATE, "r") as package:
        for src, dst in wanted.items():
            if src in package.namelist() and not dst.exists():
                dst.write_bytes(package.read(src))


def add_header(slide, title: str, number: int) -> None:
    box(slide, 0, 0, SLIDE_W, 0.82, fill=BLUE, line_color=BLUE, width=0)
    box(slide, 12.40, 0, 0.933, 0.82, fill=DARK_BLUE, line_color=DARK_BLUE, width=0)
    logo = OUT_DIR / "template_media" / "image2.png"
    if logo.exists():
        add_picture(slide, logo, 0.27, 0.15, 1.95)
    text_box(slide, title, 2.60, 0.16, 8.10, 0.42, size=35, color=WHITE, align=PP_ALIGN.CENTER, margin=0)
    text_box(slide, f"{number}/6", 12.47, 0.20, 0.78, 0.32, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0)
    box(slide, 0, 7.16, SLIDE_W, 0.34, fill=BLUE, line_color=BLUE, width=0)


def add_footer_note(slide, text: str) -> None:
    text_box(slide, text, 0.70, 6.58, 11.90, 0.30, size=24, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, margin=0)


def section_title(slide, text: str, x: float, y: float, w: float) -> None:
    text_box(slide, text, x, y, w, 0.32, size=24, bold=True, color=DARK_BLUE, margin=0)
    line(slide, x, y + 0.38, x + w, y + 0.38, DARK_BLUE, 1.2)


def bullet_list(slide, items: list[str], x: float, y: float, w: float, h: float, size: float = 23, gap: float = 0.58):
    for i, item in enumerate(items):
        text_box(slide, "•", x, y + i * gap, 0.22, 0.25, size=size, bold=True, color=DARK_BLUE, margin=0)
        text_box(slide, item, x + 0.35, y + i * gap, w - 0.35, 0.45, size=size, margin=0)


def numbered_list(slide, items: list[str], x: float, y: float, w: float, size: float = 23, gap: float = 0.62):
    for i, item in enumerate(items, 1):
        text_box(slide, f"{i}.", x, y + (i - 1) * gap, 0.35, 0.28, size=size, bold=True, color=DARK_BLUE, margin=0)
        text_box(slide, item, x + 0.45, y + (i - 1) * gap, w - 0.45, 0.48, size=size, margin=0)


def render_formula(path: Path, lines: list[str], fontsize: float, fig_w: float, fig_h: float | None = None) -> None:
    if fig_h is None:
        fig_h = 0.70 if len(lines) == 1 else 0.74 * len(lines)
    fig = plt.figure(figsize=(fig_w, fig_h), dpi=420)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ys = [0.50] if len(lines) == 1 else np.linspace(0.82, 0.18, len(lines))
    for y, s in zip(ys, lines):
        ax.text(0.5, y, s, ha="center", va="center", fontsize=fontsize, color="black")
    fig.savefig(path, dpi=420, transparent=True, bbox_inches="tight", pad_inches=0.025)
    plt.close(fig)


def generate_schematic() -> None:
    GEN.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update({"font.family": "DejaVu Serif", "mathtext.fontset": "dejavuserif"})
    fig, ax = plt.subplots(figsize=(5.8, 4.0), dpi=260)
    ax.set_aspect("equal")
    a, b, n = 6.2, 3.7, 3.0
    xs = np.arange(-8, 9)
    ys = np.arange(-6, 7)
    xx, yy = np.meshgrid(xs, ys)
    inside = (np.abs(xx / a) ** n + np.abs(yy / b) ** n) <= 1
    for x in xs:
        ax.plot([x, x], [ys.min(), ys.max()], color="#CAD6E5", lw=0.7, ls="--", zorder=0)
    for y in ys:
        ax.plot([xs.min(), xs.max()], [y, y], color="#CAD6E5", lw=0.7, ls="--", zorder=0)
    ax.scatter(xx[~inside], yy[~inside], s=20, color="white", edgecolors="#1E5AA5", linewidths=0.8, zorder=2)
    ax.scatter(xx[inside], yy[inside], s=36, color="#1E5AA5", edgecolors="white", linewidths=0.5, zorder=3)
    t = np.linspace(0, 2 * np.pi, 800)
    bx = np.sign(np.cos(t)) * np.abs(np.cos(t)) ** (2 / n) * a
    by = np.sign(np.sin(t)) * np.abs(np.sin(t)) ** (2 / n) * b
    ax.plot(bx, by, color="black", lw=2.0, zorder=4)
    ax.axis("off")
    fig.tight_layout(pad=0.02)
    fig.savefig(GEN / "lattice_superellipse.png", bbox_inches="tight")
    plt.close(fig)


def generate_formula_assets() -> None:
    GEN.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update({"font.family": "DejaVu Serif", "mathtext.fontset": "dejavuserif"})
    render_formula(
        GEN / "f_hamiltonian.png",
        [
            r"$H=\sum_i \varepsilon_i |i\rangle\langle i|$",
            r"$+\sum_{\langle i,j\rangle}t_{ij}\left(|i\rangle\langle j|+|j\rangle\langle i|\right),$",
        ],
        fontsize=24,
        fig_w=6.6,
        fig_h=1.25,
    )
    render_formula(GEN / "f_params.png", [r"$\varepsilon_i=0,\qquad t_{ij}=-1.$"], fontsize=27, fig_w=4.0)
    render_formula(GEN / "f_ekin.png", [r"$E_{\mathrm{kin}}=E_0+4.$"], fontsize=31, fig_w=3.6)
    render_formula(
        GEN / "f_superellipse.png",
        [r"$\left|\frac{x}{a}\right|^n+\left|\frac{y}{b}\right|^n\leq 1,\qquad b=a r_{\mathrm{AR}}.$"],
        fontsize=28,
        fig_w=6.8,
    )
    render_formula(
        GEN / "f_features.png",
        [r"$\left[\frac{1}{a^2},\quad\frac{1}{a^2 r_{\mathrm{AR}}},\quad r_{\mathrm{AR}}\right]$"],
        fontsize=30,
        fig_w=5.5,
    )


def add_validation_grid(slide, x: float, y: float, title: str, mode: str) -> None:
    text_box(slide, title, x, y, 2.35, 0.28, size=21, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, margin=0)
    py = y + 0.44
    cell = 0.21
    gap = 0.025
    for r in range(5):
        for c in range(5):
            active = (mode == "a" and c == 2) or (mode == "r" and r == 2)
            fill = BLUE if active else LIGHT_BLUE
            box(slide, x + 0.52 + c * (cell + gap), py + r * (cell + gap), cell, cell, fill=fill, line_color=WHITE, width=0.6)
    label = "исключение одного значения a" if mode == "a" else "исключение одного отношения сторон"
    text_box(slide, label, x - 0.06, py + 1.28, 2.45, 0.42, size=17, align=PP_ALIGN.CENTER, margin=0)


def build() -> Presentation:
    ensure_template_logos()
    generate_schematic()
    generate_formula_assets()

    if TEMPLATE.exists():
        prs = Presentation(str(TEMPLATE))
        delete_all_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = inch(SLIDE_W)
        prs.slide_height = inch(SLIDE_H)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE, line_color=BLUE, width=0)
    media = OUT_DIR / "template_media"
    for path, x, y, w in [
        (media / "image2.png", 0.55, 0.34, 2.75),
        (media / "image5.png", 3.70, 0.40, 2.60),
        (media / "image6.png", 6.85, 0.36, 3.10),
    ]:
        if path.exists():
            add_picture(slide, path, x, y, w)
    text_box(slide, "ДИПЛОМНАЯ РАБОТА", 0.80, 1.82, 11.70, 0.38, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text_box(
        slide,
        "Моделирование энергетического спектра\nи волновых функций в квантовых точках\nсложной геометрии",
        0.95,
        2.34,
        11.50,
        1.35,
        size=32,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    text_box(
        slide,
        "Модельные суперэллиптические квантовые точки\nна квадратной решётке метода сильной связи",
        1.15,
        4.06,
        11.05,
        0.70,
        size=25,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    line(slide, 1.20, 4.96, 12.10, 4.96, WHITE, 1.0)
    text_box(
        slide,
        "Докладчик: Альвинский Александр Андреевич, студент 5 курса\n"
        "Научный руководитель: А. В. Ларькин, канд. физ.-мат. наук\n"
        "Минск, 2026",
        1.00,
        5.34,
        11.35,
        0.95,
        size=22,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    text_box(slide, "1/6", 12.45, 0.20, 0.82, 0.34, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margin=0)

    # 2. Relevance, aim and tasks
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Актуальность, цель и задачи", 2)
    text_box(
        slide,
        "Форма и размер модельной квантовой точки управляют низкоэнергетическим спектром за счёт квантового ограничения.",
        0.75,
        1.20,
        11.95,
        0.58,
        size=25,
        bold=True,
        color=DARK_BLUE,
        align=PP_ALIGN.CENTER,
    )
    section_title(slide, "Цель", 0.75, 2.25, 5.30)
    text_box(
        slide,
        "Смоделировать спектр и волновые функции квантовых точек сложной геометрии и проверить физически мотивированную аппроксимацию низкоэнергетических характеристик.",
        0.75,
        2.90,
        5.35,
        1.65,
        size=23,
    )
    section_title(slide, "Задачи", 6.85, 2.25, 5.35)
    numbered_list(
        slide,
        [
            "Построить модель метода сильной связи.",
            "Проверить физическую согласованность расчётов.",
            "Сравнить Ridge и MLP при структурированной проверке.",
        ],
        6.85,
        2.90,
        5.25,
        size=23,
        gap=0.72,
    )
    add_footer_note(slide, "Сначала физическая проверка, затем суррогатная аппроксимация.")

    # 3. Method
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Метод исследования", 3)
    add_picture(slide, GEN / "f_hamiltonian.png", 0.55, 1.15, 5.65)
    text_box(slide, "(1)", 5.85, 1.58, 0.40, 0.22, size=20, color=BLACK, align=PP_ALIGN.CENTER, margin=0)
    add_picture(slide, GEN / "f_params.png", 1.55, 2.22, 3.55)
    text_box(
        slide,
        "где εᵢ — энергия узла,\ntᵢⱼ — интеграл перескока между ближайшими узлами.",
        0.82,
        2.92,
        5.20,
        0.65,
        size=18,
        color=BLACK,
    )
    add_picture(slide, GEN / "f_ekin.png", 1.63, 3.78, 3.10)
    text_box(slide, "(2)", 5.05, 4.04, 0.40, 0.22, size=20, color=BLACK, align=PP_ALIGN.CENTER, margin=0)
    text_box(slide, "где −4 — дно зоны бесконечной квадратной решётки.", 0.82, 4.48, 5.35, 0.32, size=18)
    add_picture(slide, GEN / "f_superellipse.png", 0.55, 5.00, 5.75)
    text_box(slide, "(3)", 5.95, 5.20, 0.40, 0.22, size=20, color=BLACK, align=PP_ALIGN.CENTER, margin=0)
    text_box(slide, "где a — размер, r_AR — отношение сторон, n — класс формы.", 0.82, 5.72, 5.35, 0.32, size=18)
    add_picture(slide, GEN / "lattice_superellipse.png", 6.95, 1.25, 5.35, 4.20)
    text_box(
        slide,
        "Рисунок 1 — Квадратная решётка и суперэллиптическая область",
        6.75,
        5.60,
        5.80,
        0.42,
        size=18,
        align=PP_ALIGN.CENTER,
    )

    # 4. Verification
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Верификация расчётов", 4)
    section_title(slide, "Физические проверки", 0.72, 1.25, 4.45)
    bullet_list(
        slide,
        [
            "прямоугольная контрольная задача;",
            "проверка масштаба 1/a²;",
            "круговая проверка Бесселя: ошибка ≈ 2.03%;",
            "проверка при фиксированной форме: максимум отклонения ≈ 2.12%.",
        ],
        0.75,
        2.00,
        4.45,
        2.70,
        size=22,
        gap=0.70,
    )
    add_picture(slide, ASSETS / "ar_scaling_relative_deviation.png", 5.25, 1.18, 7.20, 4.65)
    text_box(
        slide,
        "Рисунок 2 — Максимальное относительное отклонение величины (E₀ + 4)a² от среднего значения при фиксированной форме",
        5.20,
        5.95,
        7.35,
        0.52,
        size=18,
        align=PP_ALIGN.CENTER,
    )

    # 5. Surrogate results
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Суррогатная аппроксимация", 5)
    section_title(slide, "Физически мотивированные признаки", 0.70, 1.15, 5.75)
    add_picture(slide, GEN / "f_features.png", 1.42, 1.75, 4.05)
    text_box(
        slide,
        "первый признак задаёт масштаб квантового ограничения;\n"
        "второй — площадь-подобную поправку;\n"
        "третий — анизотропию формы.",
        0.78,
        2.58,
        5.55,
        1.10,
        size=19,
    )
    text_box(
        slide,
        "Ridge-регрессия — основная интерпретируемая модель.\n"
        "MLP — многослойный перцептрон, контрольная нелинейная модель.\n"
        "Случайное разбиение не использовалось.",
        0.78,
        3.78,
        5.80,
        1.22,
        size=20,
        margin=0,
    )
    section_title(slide, "Структурированная проверка", 7.00, 1.15, 5.40)
    add_validation_grid(slide, 7.15, 1.72, "LOAO", "a")
    add_validation_grid(slide, 9.80, 1.72, "LOARO", "r")
    text_box(slide, "MLP + физические признаки", 7.05, 3.96, 5.20, 0.28, size=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, margin=0)
    text_box(slide, "10/16 ячеек\nLOAO: 3/8\nLOARO: 7/8", 7.30, 4.36, 4.75, 1.00, size=25, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_footer_note(slide, "MLP не показал устойчивого преимущества по заданному критерию.")

    # 6. Conclusions and limitations
    slide = prs.slides.add_slide(blank)
    add_header(slide, "Выводы и ограничения", 6)
    section_title(slide, "Основные выводы", 0.70, 1.10, 6.15)
    numbered_list(
        slide,
        [
            "Спектр согласуется с масштабом квантового ограничения.",
            "Ridge-регрессия остаётся основной интерпретируемой моделью.",
            "MLP не дал устойчивого преимущества.",
            "Ошибка Ridge: E₀ 2.01–2.60% от Eₖᵢₙ; dE₁ 0.78–1.15%.",
            "Простые диагностики границы не объясняют остатки универсально.",
        ],
        0.72,
        1.82,
        6.30,
        size=20,
        gap=0.62,
    )
    section_title(slide, "Ограничения", 7.30, 1.10, 5.20)
    bullet_list(
        slide,
        [
            "только заданная сетка параметров;",
            "n рассматривается как дискретный класс формы;",
            "без DFT/OpenMX-калибровки;",
            "суррогат не заменяет прямой расчёт Kwant.",
        ],
        7.32,
        1.85,
        5.15,
        2.35,
        size=21,
        gap=0.62,
    )
    text_box(
        slide,
        "Прямой расчёт Kwant остаётся эталоном.",
        7.45,
        4.85,
        4.90,
        0.70,
        size=26,
        bold=True,
        color=DARK_BLUE,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    line(slide, 7.55, 5.70, 12.10, 5.70, DARK_BLUE, 1.3)

    return prs


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(OUT)
    print(OUT)
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
